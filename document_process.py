import os
import re
import shutil
import traceback
from collections import Counter
from typing import List
from urllib.parse import quote, unquote
import uvicorn
import openpyxl
from bs4 import BeautifulSoup
from docx import Document
from docx.table import _Cell
from fastapi import FastAPI, File, UploadFile, Query, Form
from fastapi.encoders import jsonable_encoder
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import JSONResponse, FileResponse
from pptx import Presentation
from apscheduler.schedulers.background import BackgroundScheduler
from pathlib import Path
import requests
import threading
from CustomThread import CustomThread
from TaskCounter import TaskCounter


app = FastAPI()

app.add_middleware(CORSMiddleware, allow_origins=["*"],
    allow_credentials=False, allow_methods=["*"],
    allow_headers=["*"],
)

TRANSLATOR_URL = "http://dilato-webtrans-nmt-patent:7878/translation"
req_storage_path = "./req_doc"
storage_path = "./translated"
src_lang_list = ["ko", "en"]
tgt_lang_list = ["en", "ko"]
valid_lang_list = ["ko", "en"]

def clean_folders():
    try:
        input_dir = Path("req_doc")
        output_dir = Path("translated")

        for dir_path in [input_dir, output_dir]:
            if dir_path.exists():
                shutil.rmtree(dir_path)
            dir_path.mkdir()

    except Exception as e:
        pass

# 한국시간 새벽 4시에 폴더를 초기화한다.
scheduler = BackgroundScheduler(daemon=True, timezone='Asia/Seoul')
scheduler.add_job(clean_folders, 'cron', hour=4, minute=0)
scheduler.start()


# [번역전처리] 문장의 non-character 여부 검사
def is_character(sentence):
    pattern = r'[가-힣a-zA-Z0-9]'
    return bool(re.search(pattern, sentence))


# [번역] source language => target language 번역 API
def translate(text, src, tgt):
    headers = {"Content-Type": "application/json"}
    data = {"text": text, "from_lang": src, "to_lang": tgt}

    try:
        session = requests.Session()
        response = session.post(
            TRANSLATOR_URL,
            headers=headers,
            json=data,
        )
        return response.json()['translation']['result']
    except Exception as e:
        print(f"번역요청에서 에러 발생 : [{e}]")
        return ""

def extract_text_with_ids(soup):
    text_elements = []
    element_seq = 1

    for element in soup.find_all(True):  # True means all tags
        if element.name not in ['style', 'meta', 'script']:
            if element.string and element.string.strip():  # Check if element has text
                # Assigning a unique ID
                element_id = f'element_{element_seq}'
                text_elements.append({'id': element_id, 'original_text': element.string.strip(),
                    'translated_text': element.string.strip(), 'tag': element.name, 'parent': str(element.parent.name)})
                element.string.replace_with(element_id)
                element_seq += 1
    return text_elements


def replace_html_text(soup, sentences):
    for sentence in sentences:
        element = soup.find(True, string=sentence["id"])
        if element and element.string and element.string.strip():
            element.string.replace_with(sentence["translated_text"])
    return soup

def parse_html(html_text, from_lang, to_lang):
    detection_list = []
    # BeautifulSoup을 사용하여 HTML 파싱
    html_text = html_text.replace("<br/>", " ")
    html_text = html_text.replace("<br />", " ")
    html_text = html_text.replace("<br>", " ")
    soup = BeautifulSoup(html_text, 'html.parser')
    soup = BeautifulSoup(soup.prettify(), 'html.parser')

    sentences = extract_text_with_ids(soup)
    translation_results = []
    for sentence in sentences:
        translation_results.append(translate(sentence, from_lang, to_lang))

    if translation_results:
        for i, result in enumerate(translation_results):
            sentences[i]["translated_text"] = result["translations"]
            detection_list.append(result["detect_lang"])
        soup = replace_html_text(soup, sentences)
        translated_html_text = str(soup)
        return detection_list, translated_html_text
    else:
        detection_list = ["nn"]

    return detection_list, html_text

# [HTML]
def html_translation(input_html, output_html, uuid, from_lang="ko", to_lang="en"):
    uuid_folder_path = os.path.join(req_storage_path, uuid)
    os.makedirs(uuid_folder_path, exist_ok=True)

    with open(os.path.join(uuid_folder_path, input_html), 'r', encoding='utf-8') as file:
        input_html_data = file.read()

    detection_list, output_html_data = parse_html(input_html_data, from_lang, to_lang)
    print("[DL] : ", detection_list)
    counter = Counter(detection_list)
    print(counter)
    most_common = counter.most_common(1)
    detect_lang = most_common[0][0]

    with open(os.path.join(uuid_folder_path, "html", output_html), 'w', encoding='utf-8') as file:
        file.write(output_html_data)


def only_allowed_chars(text):
    # 허용 문자: 숫자, _, -, *, %, 공백
    pattern = r'^[0-9_\-\*\% ]+$'
    return bool(re.match(pattern, text))

# [DOCX]
def docx_translation(input_docx, output_docx, uuid, from_lang="ko", to_lang="en"):
    try:
        uuid_folder_path = os.path.join(req_storage_path, uuid)
        os.makedirs(uuid_folder_path, exist_ok=True)
        uuid_storage_path = os.path.join(storage_path, uuid)
        os.makedirs(uuid_storage_path, exist_ok=True)

        docx = Document(os.path.join(uuid_folder_path, input_docx))

        for para in docx.paragraphs:
            if TaskCounter.task_dict[uuid]["task"].stop_event.is_set():
                raise SystemExit("스레드를 종료합니다.")

            original_text = para.text.strip()

            if original_text == "" or only_allowed_chars(original_text):
                continue

            para.text = translate(original_text, from_lang, to_lang)

        for table in docx.tables:
            for row in table._tbl.tr_lst:  # lxml element 순회
                for tc in row.tc_lst:
                    if TaskCounter.task_dict[uuid]["task"].stop_event.is_set():
                        raise SystemExit("스레드를 종료합니다.")

                    # 병합된 셀 중복 출력 방지
                    if tc.vMerge == 'continue':
                        continue
                    cell = _Cell(tc, table)
                    text = cell.text.strip()

                    if text == "" or only_allowed_chars(text) :
                        continue

                    text = translate(cell.text, from_lang, to_lang)
                    cell.text = text

        docx.save(os.path.join(uuid_storage_path, output_docx))

    except SystemExit as e:
        print("번역이 취소되었습니다.")
    except Exception as e:
        print("문서번역 중 예상치 못한 에러가 발생했습니다.")
    finally:
        del TaskCounter.task_dict[uuid]

# [PPTX]
def pptx_translation(input_pptx, output_pptx, uuid, from_lang="ko", to_lang="en"):
    try:
        uuid_folder_path = os.path.join(req_storage_path, uuid)
        os.makedirs(uuid_folder_path, exist_ok=True)
        uuid_storage_path = os.path.join(storage_path, uuid)
        os.makedirs(uuid_storage_path, exist_ok=True)

        prs = Presentation(os.path.join(uuid_folder_path, input_pptx))

        for slide in prs.slides:
            for shape in slide.shapes:
                if TaskCounter.task_dict[uuid]["task"].stop_event.is_set():
                    raise SystemExit("스레드를 종료합니다.")

                # 텍스트를 포함하는 도형 수정
                if hasattr(shape, "text"):
                    for paragraph in shape.text_frame.paragraphs:
                        if TaskCounter.task_dict[uuid]["task"].stop_event.is_set():
                            raise SystemExit("스레드를 종료합니다.")

                        translated_text = translate(paragraph.text, from_lang, to_lang)

                        if paragraph.level > 0:
                            try:
                                bullet_style = paragraph.bullet
                                paragraph.text = translated_text
                                paragraph.bullet = bullet_style
                            except:
                                paragraph.text = translated_text
                        else:
                            paragraph.text = translated_text

        prs.save(os.path.join(uuid_storage_path, output_pptx))

    except SystemExit as e:
        print("번역이 취소되었습니다.")
    except Exception as e:
        print("문서번역 중 예상치 못한 에러가 발생했습니다.")
    finally:
        del TaskCounter.task_dict[uuid]

# [XLSX]
def xlsx_translation(input_xlsx, output_xlsx, uuid, from_lang="ko", to_lang="en"):
    try:
        uuid_folder_path = os.path.join(req_storage_path, uuid)
        os.makedirs(uuid_folder_path, exist_ok=True)
        uuid_storage_path = os.path.join(storage_path, uuid)
        os.makedirs(uuid_storage_path, exist_ok=True)

        workbook = openpyxl.load_workbook(os.path.join(uuid_folder_path, input_xlsx))

        for sheet in workbook.sheetnames:
            worksheet = workbook[sheet]
            for row in worksheet.iter_rows():
                for cell in row:
                    if TaskCounter.task_dict[uuid]["task"].stop_event.is_set():
                        raise SystemExit("스레드를 종료합니다.")

                    if isinstance(cell.value, str) and cell.value.strip():
                        cell.value = translate(cell.value, from_lang, to_lang)

        workbook.save(os.path.join(uuid_storage_path, output_xlsx))

    except SystemExit as e:
        print("번역이 취소되었습니다.")
    except Exception as e:
        print("문서번역 중 예상치 못한 에러가 발생했습니다.")
    finally:
        del TaskCounter.task_dict[uuid]

# [TXT]
def txt_translation(input_txt, output_txt, uuid, from_lang="ko", to_lang="en"):
    try:
        uuid_folder_path = os.path.join(req_storage_path, uuid)
        os.makedirs(uuid_folder_path, exist_ok=True)
        uuid_storage_path = os.path.join(storage_path, uuid)
        os.makedirs(uuid_storage_path, exist_ok=True)

        with open(os.path.join(uuid_folder_path, input_txt), 'r', encoding='utf-8') as file:
            lines = file.readlines()

        for i in range(len(lines)):
            if TaskCounter.task_dict[uuid]["task"].stop_event.is_set():
                raise SystemExit("스레드를 종료합니다.")

            text = lines[i].strip()
            if text == "" or only_allowed_chars(text):
                continue

            lines[i] = translate(text, from_lang, to_lang)

        with open(os.path.join(uuid_storage_path, output_txt), 'w', encoding='utf-8') as new_file:
            new_file.writelines(lines)

    except SystemExit as e:
        print("번역이 취소되었습니다.")
    except Exception as e:
        print("문서번역 중 예상치 못한 에러가 발생했습니다.")
    finally:
        del TaskCounter.task_dict[uuid]


# [파일 번역 처리]
def run_document_translation(uuid, input_file, src="ko", tgt="en"):
    _, file_ext = os.path.splitext(input_file)
    output_filename = input_file

    if file_ext == '.html':
        html_translation(input_file, output_filename, uuid, src, tgt)
    elif file_ext == '.txt':
        txt_translation(input_file, output_filename, uuid, src, tgt)
    elif file_ext == '.docx':
        docx_translation(input_file, output_filename, uuid, src, tgt)
    elif file_ext == ".pptx":
        pptx_translation(input_file, output_filename, uuid, src, tgt)
    elif file_ext == ".xlsx":
        xlsx_translation(input_file, output_filename, uuid, src, tgt)


# [서버통신확인용]
@app.get("/ping")
async def pong():
    json_compatible_data = jsonable_encoder({'response': "pong"})
    return JSONResponse(json_compatible_data)


# [파일 번역 요청] Query문을 이용하여 사용 (format, filename)
@app.post("/trans_file")
async def trans_file(uuid: str = Form(...), from_lang: str = Form(...),
        to_lang: str = Form(...), files: List[UploadFile] = File(...)):
    try:
        for file in files:
            filename = file.filename

            uuid_folder_path = os.path.join(req_storage_path, uuid)
            req_file_path = os.path.join(uuid_folder_path, filename)
            os.makedirs(uuid_folder_path, exist_ok=True)

            uuid_storage_path = os.path.join(storage_path, uuid)
            os.makedirs(uuid_storage_path, exist_ok=True)

            # 파일 저장
            with open(req_file_path, "wb") as buffer:
                buffer.write(await file.read())

            TaskCounter.task_dict[uuid] = {"task": None, "counter": None}
            stop_event = threading.Event()
            thread = threading.Thread(target=run_document_translation, args=(uuid, filename, from_lang, to_lang), daemon = False)
            thread.start()
            TaskCounter.task_dict[uuid]["task"] = CustomThread(thread, stop_event)

        return JSONResponse(content={"task": "processing"}, status_code=200)

    except Exception as e:
        return JSONResponse(content={"message": "Internal server error"}, status_code=500)


# [번역된 파일 엔드포인트] API를 호출하여 파일 엔드포인트 접근
@app.get("/download/{file_name}")
async def download_file(file_name: str, uuid: str = Query(...)):
    try:
        file_name = unquote(file_name)
        uuid_storage_path = os.path.join(storage_path, uuid)
        file_path = os.path.join(uuid_storage_path, file_name)

        if os.path.exists(file_path):
            return FileResponse(file_path, media_type='application/octet-stream')
        else:
            return JSONResponse(content={"error": "File not found."}, status_code=404)

    except Exception as e:
        return JSONResponse(content={"message": "Internal server error"}, status_code=500)


# [파일 존재 여부 확인]
@app.get("/check_file")
async def check_file(uuid: str = Query(...), filename: str = Query(...)):
    try:
        uuid_storage_path = os.path.join(storage_path, uuid)
        target_file_name = unquote(filename)
        target_file_path = os.path.join(uuid_storage_path, target_file_name)

        if os.path.exists(target_file_path):
            target_file_size = os.path.getsize(target_file_path)
            return JSONResponse(content={"exists": True, "filename": quote(target_file_name), "size": target_file_size},
                                status_code=200)
        else:
            return JSONResponse(content={"exists": False, "path": "None"}, status_code=200)

    except Exception as e:
        return JSONResponse(content={"message": "Internal server error"}, status_code=500)

@app.delete("/cancel")
def cancel_trans(uuid: str):
    try:
        task_obj = TaskCounter.task_dict.get(uuid)
        if task_obj:
            task_obj["task"].stop()

        return JSONResponse(status_code=200, content={"message": "The task has already been completed or deleted."})

    except Exception as e:
        return JSONResponse(content={"message": "Internal server error"}, status_code=500)

if __name__ == "__main__":
    uvicorn.run(app, host="0.0.0.0", port=8080)
